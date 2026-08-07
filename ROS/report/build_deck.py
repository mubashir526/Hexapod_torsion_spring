#!/usr/bin/env python3
"""Build the STRP Check-In 2 results deck, styled to match presentation2.pdf's theme."""

import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

CODE = Path("/home/mubashir/Pictures/FYP-Legged-Robot-main/Code")
FIG = CODE / "ROS" / "report" / "figures"   # verified, data-driven figures
OUT = CODE / "STRP_CheckIn2_Simulation_Results.pptx"

# ---------------------------------------------------------------- palette
NAVY = RGBColor(0x2B, 0x50, 0x98)      # slide background
NAVY_DK = RGBColor(0x1E, 0x39, 0x6E)   # table headers
PERI = RGBColor(0x71, 0x81, 0xB8)      # periwinkle panel
CARD = RGBColor(0xDD, 0xE5, 0xFF)      # light blue card
CARD2 = RGBColor(0xED, 0xF1, 0xFF)     # alt row
CREAM = RGBColor(0xF5, 0xEF, 0xEF)     # light text
YELLOW = RGBColor(0xFF, 0xD2, 0x5E)    # accent
TEAL = RGBColor(0x64, 0xCB, 0xD6)      # accent
INK = RGBColor(0x11, 0x06, 0x66)       # dark heading on light
INK2 = RGBColor(0x16, 0x25, 0x5C)      # body on light
MUTED = RGBColor(0x9E, 0xB0, 0xDB)     # footer

SERIF = "Georgia"
SANS = "Trebuchet MS"

SW, SH = 13.333, 7.5
L, R = 0.62, 12.71          # content margins
CW = R - L                  # 12.09 usable width
NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"  # No Style, No Grid

TOTAL_SLIDES = 13


# ---------------------------------------------------------------- helpers
def solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def rect(slide, x, y, w, h, rgb, shape=MSO_SHAPE.RECTANGLE, radius=None, rot=0):
    sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    solid(sh, rgb)
    if radius is not None:
        try:
            sh.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    if rot:
        sh.rotation = rot
    return sh


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, rgb, font=SANS, bold=False, italic=False, first=False,
         space_before=0, space_after=4, align=PP_ALIGN.LEFT, bullet=None,
         line_spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line_spacing
    if bullet:
        r = p.add_run()
        r.text = bullet + "  "
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = YELLOW if rgb in (CREAM, CARD) else TEAL
        r.font.name = SANS
    for seg, is_b in _split_bold(text):
        r = p.add_run()
        r.text = seg
        r.font.size = Pt(size)
        r.font.bold = bold or is_b
        r.font.italic = italic
        r.font.color.rgb = rgb
        r.font.name = font
    return p


def _split_bold(text):
    """Split on **bold** markers -> [(segment, is_bold), ...]."""
    out, buf, i = [], "", 0
    while i < len(text):
        if text.startswith("**", i):
            j = text.find("**", i + 2)
            if j == -1:
                buf += text[i:]
                break
            if buf:
                out.append((buf, False))
                buf = ""
            out.append((text[i + 2:j], True))
            i = j + 2
        else:
            buf += text[i]
            i += 1
    if buf:
        out.append((buf, False))
    return out or [("", False)]


def zigzag(slide, x, y, seg_h=0.34, seg_w=0.26, n=4, thick=0.055, rgb=YELLOW):
    """Yellow lightning zigzag, as in the template's left edge."""
    pts = [(x + (seg_w if i % 2 else 0.0), y + i * seg_h) for i in range(n + 1)]
    fwd = [(px + thick, py) for px, py in pts]
    back = [(px - thick, py) for px, py in reversed(pts)]
    path = fwd + back
    b = slide.shapes.build_freeform(Inches(path[0][0]), Inches(path[0][1]))
    b.add_line_segments([(Inches(px), Inches(py)) for px, py in path[1:]], close=True)
    solid(b.convert_to_shape(), rgb)


def stripes(slide, x, y, length=1.62, thick=0.105, gap=0.175,
            colors=(TEAL, YELLOW, CREAM, YELLOW)):
    """Diagonal accent stripes for the bottom-left corner."""
    for i, c in enumerate(colors):
        rect(slide, x + i * gap * 0.72, y + i * gap * 0.72, length, thick, c, rot=-45)


def star(slide, x, y, size, rgb=CREAM):
    sh = rect(slide, x, y, size, size, rgb, shape=MSO_SHAPE.STAR_4_POINT)
    try:
        sh.adjustments[0] = 0.10  # slimmer points, like the template's sparkles
    except (IndexError, ValueError):
        pass


def table_no_style(tbl):
    """Strip PowerPoint's default banded table style so explicit fills read clean."""
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is None:
        return
    for tag in ("a:tableStyleId",):
        el = tblPr.find(qn(tag))
        if el is not None:
            tblPr.remove(el)
    el = tblPr.makeelement(qn("a:tableStyleId"), {})
    el.text = NO_STYLE
    tblPr.append(el)
    tblPr.set("bandRow", "0")
    tblPr.set("firstRow", "0")


def add_table(slide, x, y, w, rows, col_w=None, fs=10.0, hdr_fs=None,
              row_h=0.30, hdr_h=None, highlight=(), warn=(),
              aligns=None, hdr_rgb=NAVY_DK):
    """rows[0] is the header. col_w = list of relative widths."""
    nr, nc = len(rows), len(rows[0])
    hdr_fs = hdr_fs or fs
    hdr_h = hdr_h or row_h + 0.04
    h = hdr_h + (nr - 1) * row_h
    gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = gf.table
    table_no_style(tbl)
    tbl.first_row = True
    tbl.horz_banding = False

    if col_w:
        tot = float(sum(col_w))
        for i, cwv in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * cwv / tot))
    tbl.rows[0].height = Inches(hdr_h)
    for i in range(1, nr):
        tbl.rows[i].height = Inches(row_h)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.055)
            cell.margin_top = cell.margin_bottom = Inches(0.014)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                bg, fg, bold, size = hdr_rgb, CREAM, True, hdr_fs
            elif ri in highlight:
                bg, fg, bold, size = YELLOW, INK, True, fs
            elif ri in warn:
                bg, fg, bold, size = RGBColor(0xF2, 0xC8, 0xC0), INK2, False, fs
            else:
                bg, fg, bold, size = (CARD if ri % 2 else CARD2), INK2, False, fs
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = (aligns[ci] if aligns else
                           (PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER))
            for seg, is_b in _split_bold(str(val)):
                r = p.add_run()
                r.text = seg
                r.font.size = Pt(size)
                r.font.bold = bold or is_b
                r.font.color.rgb = fg
                r.font.name = SANS
    return gf, h


def add_pic(slide, path, x, y, w, h, center=True, frame=True):
    """Fit image inside (x,y,w,h) preserving aspect; optional cream mat behind it."""
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w / h > ar:
        ph, pw = h, h * ar
    else:
        pw, ph = w, w / ar
    px = x + (w - pw) / 2 if center else x
    py = y + (h - ph) / 2 if center else y
    if frame:
        pad = 0.07
        rect(slide, px - pad, py - pad, pw + 2 * pad, ph + 2 * pad, CREAM,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.035)
    slide.shapes.add_picture(str(path), Inches(px), Inches(py),
                             Inches(pw), Inches(ph))
    return px, py, pw, ph


# ---------------------------------------------------------------- slide frame
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]
_n = [0]


def new_slide(title=None, caption=None, deco="content"):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)

    if deco == "title":
        b = s.shapes.build_freeform(Inches(0), Inches(SH))
        b.add_line_segments([(Inches(0), Inches(1.55)), (Inches(SW), Inches(SH))],
                            close=True)
        solid(b.convert_to_shape(), PERI)
        star(s, 8.55, 0.42, 0.62)
        star(s, 11.72, 1.28, 0.70)
        star(s, 5.55, 6.48, 0.40)
        star(s, 1.62, 1.15, 0.34)
    else:
        zigzag(s, 0.11, 1.72)
        stripes(s, 0.02, 7.16)
        star(s, 12.76, 0.30, 0.40)
        star(s, 0.13, 6.60, 0.28)

    _n[0] += 1
    if title is not None:
        ts = 27 if len(title) <= 44 else (24 if len(title) <= 54 else 21)
        tf = textbox(s, L, 0.26, CW - 0.35, 0.66)
        para(tf, title, ts, CREAM, font=SERIF, bold=True, first=True, space_after=0)
        rect(s, L, 0.99, 1.75, 0.045, YELLOW)
    if caption is not None:
        tf = textbox(s, L, 1.10, CW, 0.40)
        para(tf, caption, 11.5, CARD, first=True, space_after=0, line_spacing=1.05)

    if deco != "title":
        tf = textbox(s, R - 4.6, 7.03, 4.6, 0.28)
        para(tf, f"STRP Check-In 2  ·  Simulation & Results  ·  {_n[0]} / {TOTAL_SLIDES}",
             8.5, MUTED, first=True, space_after=0, align=PP_ALIGN.RIGHT)
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


TOP_C = 1.56   # content top when a caption is present
TOP_N = 1.24   # content top with no caption
BOT = 6.88


# ================================================================ 1. TITLE
s = new_slide(deco="title")
rect(s, 2.35, 1.32, 8.65, 4.42, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
tf = textbox(s, 2.75, 1.72, 7.85, 3.6, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Passive Gravity Compensation", 36, INK, font=SERIF, bold=True,
     first=True, align=PP_ALIGN.CENTER, space_after=2)
para(tf, "for", 19, RGBColor(0x2B, 0x1A, 0x8C), font=SERIF, bold=True,
     align=PP_ALIGN.CENTER, space_after=2)
para(tf, "Graceful Degradation", 32, INK, font=SERIF, bold=True,
     align=PP_ALIGN.CENTER, space_after=2)
para(tf, "of", 19, RGBColor(0x2B, 0x1A, 0x8C), font=SERIF, bold=True,
     align=PP_ALIGN.CENTER, space_after=2)
para(tf, "Hexapod Robots", 32, INK, font=SERIF, bold=True,
     align=PP_ALIGN.CENTER, space_after=12)
para(tf, "STRP Check-In 2   |   Simulation & Experimental Results",
     13, INK2, align=PP_ALIGN.CENTER, space_after=0)

tf = textbox(s, 0.72, 6.42, 4.0, 0.85)
para(tf, "PRINCIPLE INVESTIGATOR", 11, INK, bold=True, first=True, space_after=2)
para(tf, "TABSHEER ALI ASKARI", 12, RGBColor(0x1A, 0x1A, 0x1A), space_after=0)
tf = textbox(s, 6.9, 6.42, 5.75, 0.85)
para(tf, "PRESENTED BY", 11, INK, bold=True, first=True, space_after=2,
     align=PP_ALIGN.RIGHT)
para(tf, "ABEER ANSARI     MUBASHIR BAIG     RAUF ABDULLAH", 11.5,
     RGBColor(0x33, 0x33, 0x33), space_after=0, align=PP_ALIGN.RIGHT)
notes(s, """
[0:00-0:15] Since the last check-in the work has been almost entirely in
simulation: we integrated the torsion spring, built the measurement chain to
see what it does to the motors, and ran 214 runs across five phases. This is
the results talk.
""")


# ================================================================ 2. SIM CHANGES
s = new_slide("Since Check-In 1 — Cleaner Measurements",
              "Four changes to the simulator before any spring result could be trusted.")
CH2 = 4.18
rect(s, L, TOP_C, 7.35, CH2, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
tf = textbox(s, L + 0.30, TOP_C + 0.24, 6.75, CH2 - 0.45)
items = [
    ("Locked to simulated time", "The gait loop and the torque logger both run off "
     "/clock, so command and measurement are sampled on the same timebase — no "
     "wall-clock jitter in the logs."),
    ("Removed the start-up stall", "The run now begins from a pose already inside "
     "the gait cycle instead of a dead stop. That removed the start-up transient: "
     "smaller opening torque peak, less noise in the first cycle, and lower "
     "tracking error, because the controller never has to catch up from rest."),
    ("Commanded-effort plugin", "We log the effort the motor is actually commanded "
     "to produce, per joint — not the total joint wrench. So a \"34% reduction\" is "
     "34% off what the motor sees, which is the number that matters for sizing "
     "and for battery life."),
    ("Effort-vs-angle + recorded video", "All 12 joints logged at 50 Hz against "
     "joint angle, plus a video of every run. A torque peak can now be traced to "
     "a timestamp, a joint angle, and a body pose."),
]
for i, (hd, body) in enumerate(items):
    para(tf, hd, 13.5, INK, bold=True, first=(i == 0), space_before=0 if i == 0 else 9,
         space_after=2, bullet="▸")
    para(tf, body, 11, INK2, space_after=0, line_spacing=1.03)

rect(s, 8.35, TOP_C, R - 8.35, CH2, PERI, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.05)
tf = textbox(s, 8.62, TOP_C + 0.22, 3.85, 0.34)
para(tf, "SIMULATION SETUP", 11.5, CREAM, bold=True, first=True, space_after=0)
add_table(s, 8.62, TOP_C + 0.62, 3.85, [
    ["Parameter", "Value"],
    ["Platform", "THex Quadruped, 4 × 3 DoF"],
    ["Mass m", "1.39847 kg (13 links)"],
    ["Actuator limit", "±0.9414 N·m / joint"],
    ["Control", "Position PID, 10 Hz gait"],
    ["Logging", "50 Hz, all 12 joints"],
    ["Simulator", "Gazebo Harmonic 8.14 (DART)"],
    ["Middleware", "ROS 2 Humble"],
    ["Gravity", "9.8 m/s² (CoT); 9.78 IMU"],
], col_w=[1.0, 1.35], fs=9.5, hdr_fs=9.5, row_h=0.355)

rect(s, L, TOP_C + CH2 + 0.22, CW, 0.90, PERI,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
tf = textbox(s, L + 0.30, TOP_C + CH2 + 0.36, CW - 0.60, 0.62)
para(tf, "Net effect: every number in the rest of this deck is **motor-commanded "
         "effort**, on a jitter-free timebase, with no start-up transient "
         "contaminating the first gait cycle.", 11.5, CREAM, first=True,
     space_after=0, bullet="▸", line_spacing=1.03)
notes(s, """
[0:15-0:38] Four infrastructure changes first. Sim-time locking removed log
jitter. Starting from inside the gait cycle instead of a stall removed the
opening torque spike and dropped tracking error. The commanded-effort plugin
matters most: we now measure what the motor is asked for, not the total joint
wrench, so every reduction percentage is a real motor saving.
""")


# ================================================================ 3. ROADMAP
s = new_slide("Experiment Roadmap — 214 Simulation Runs",
              "Five phases, 12 July → 30 July. Each phase fixed a measurement problem "
              "found in the one before it.")
px, py, pw, ph = add_pic(s, FIG / "timeline.png",
                         L, TOP_C, CW, 2.05)
add_table(s, L, py + ph + 0.30, CW, [
    ["Phase", "Directory", "Runs", "What it added", "Date"],
    ["1 — Baseline", "experiment_old/", "6", "First working gait; no commanded-effort logging, no settle phase", "12 Jul"],
    ["2a — Shared sweep", "experiment_before symeetry/", "111", "Spring integrated; commanded-effort logging; kx × θ₀ grid, one θ₀ for all knees", "27 Jul"],
    ["2b — Mirrored sweep", "experiment_new/", "91", "±θ₀ mirrored per leg; body state, forward displacement, cost of transport", "30 Jul"],
    ["3a — Speed (frequency)", "experiment_speed_freq/", "3", "target_freq = 5 / 10 / 20 Hz — same trajectory, different replay speed", "30 Jul"],
    ["3b — Speed (resolution)", "experiment_speed_steps/", "3", "NUM_DATA_POINTS = 8 / 16 / 32 — same speed, different waypoint count", "30 Jul"],
    ["", "**Total**", "**214**", "", ""],
], col_w=[1.55, 1.9, 0.55, 5.1, 0.72], fs=9.5, hdr_fs=10, row_h=0.40,
    highlight=(6,), aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER,
                            PP_ALIGN.LEFT, PP_ALIGN.CENTER])
notes(s, """
[0:38-0:55] 214 runs. Phase 1 baseline gait; Phase 2a the first spring sweep,
111 runs; Phase 2b re-ran it mirrored after we found a symmetry bug, 91 runs,
and added cost of transport; Phase 3 was two speed levers. Everything from 2a
on has commanded-effort logging, so 2a and 2b are directly comparable.
""")


# ================================================================ 4. SWEEP 1
s = new_slide("Sweep 1 — One Shared Rest Angle (111 runs)",
              "A linear torsion spring in parallel with each knee:  "
              "τ_motor = τ_required − kx · (θ₀ − θ).  Two knobs — stiffness kx, rest "
              "angle θ₀.")
CH4 = 4.32
rect(s, L, TOP_C, 6.05, CH4, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
tf = textbox(s, L + 0.28, TOP_C + 0.24, 5.5, 1.2)
para(tf, "Grid: kx = 0.05 → 0.50 N·m/rad (10 values) × θ₀ = 0° → −50° (11 values) "
         "= 110 spring cells + 1 baseline.", 11.5, INK2, first=True, space_after=9,
     bullet="▸", line_spacing=1.03)
para(tf, "The same θ₀ was applied to all four knees.", 11.5, INK, bold=True,
     space_after=0, bullet="▸")
tf = textbox(s, L + 0.28, TOP_C + 1.52, 5.5, 0.34)
para(tf, "TOP 5 OF 110 CELLS   ·   baseline mean knee effort 0.2345 N·m",
     9.5, INK, bold=True, first=True, space_after=0)
add_table(s, L + 0.28, TOP_C + 1.90, 5.5, [
    ["#", "kx", "θ₀", "Mean effort", "Reduction", "Track err"],
    ["1", "0.30", "0°", "0.1548", "**34.0%**", "3.51°"],
    ["2", "0.35", "0°", "0.1556", "33.6%", "3.50°"],
    ["3", "0.30", "−5°", "0.1571", "33.0%", "3.41°"],
    ["4", "0.35", "−5°", "0.1595", "32.0%", "3.48°"],
    ["5", "0.25", "0°", "0.1599", "31.8%", "3.53°"],
], col_w=[0.34, 0.6, 0.6, 1.05, 0.95, 0.9], fs=10.5, row_h=0.345, highlight=(1,),
    aligns=[PP_ALIGN.CENTER] * 6)
add_pic(s, FIG / "p2a_grids.png",
        7.05, TOP_C, R - 7.05, CH4)
rect(s, L, TOP_C + CH4 + 0.20, CW, 0.80, PERI,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.11)
tf = textbox(s, L + 0.30, TOP_C + CH4 + 0.32, CW - 0.60, 0.58)
para(tf, "Every one of the top five sits at θ₀ = 0° or −5°, and the far side of the "
         "grid collapses to −89%. That is not a coincidence — it is the problem.",
     11.5, YELLOW, bold=True, first=True, space_after=0, bullet="▸",
     line_spacing=1.03)
notes(s, """
[0:55-1:22] The spring model is one line: motor torque equals required torque
minus spring torque. Two parameters, so we swept both - 111 runs. Best cell
takes 34% off mean knee effort, 0.2345 down to 0.1548 newton-metres. But the
top five all cluster at zero or minus five degrees, and the far side of the
grid collapses to minus 89%. Something was wrong with how we applied the rest
angle.
""")


# ================================================================ 5. SYMMETRY
s = new_slide("The Symmetry Problem",
              "The legs are mirrored, but we were giving all four knees the same rest "
              "angle — so Sweep 1 was handicapped before it started.")
rect(s, L, TOP_C, 5.55, 3.62, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.05)
tf = textbox(s, L + 0.27, TOP_C + 0.24, 5.0, 0.36)
para(tf, "MEASURED STANCE OPERATING POINT", 10, INK, bold=True, first=True,
     space_after=0)
add_table(s, L + 0.27, TOP_C + 0.62, 5.0, [
    ["Knee", "Stance angle q_op", "Holding torque", "Side"],
    ["FR", "+37.2°", "−0.246 N·m", "right"],
    ["BR", "+42.9°", "−0.248 N·m", "right"],
    ["FL", "−38.4°", "+0.258 N·m", "left"],
    ["BL", "−40.8°", "+0.264 N·m", "left"],
], col_w=[0.6, 1.35, 1.35, 0.7], fs=10, row_h=0.315,
    aligns=[PP_ALIGN.CENTER] * 4)
tf = textbox(s, L + 0.27, TOP_C + 2.30, 5.0, 1.5)
para(tf, "Left and right knees sit at opposite angles and need holding torques of "
         "opposite sign. A single shared θ₀ can only serve both sides at θ₀ = 0°, "
         "where the lever arm is just −q_op — which happens to carry the right sign "
         "for each side.", 10.5, INK2, first=True, space_after=7, line_spacing=1.03)
para(tf, "Move away from 0° and the two sides fail in opposite ways.", 11, INK,
     bold=True, space_after=0)

rect(s, L, 5.42, 5.55, 1.46, PERI, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
tf = textbox(s, L + 0.28, 5.58, 5.0, 1.2)
para(tf, "So Sweep 1's \"optimum\" at θ₀ = 0° was never a discovered optimum — it "
         "was the **only rest angle a shared-θ₀ design could use**. The sweep was "
         "answering a question we had accidentally constrained.", 11, CREAM,
     first=True, space_after=0, bullet="▸", line_spacing=1.04)

rect(s, 6.42, TOP_C, R - 6.42, BOT - TOP_C, PERI,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
tf = textbox(s, 6.70, TOP_C + 0.22, 5.75, 0.34)
para(tf, "TWO DISTINCT FAILURE MODES", 11.5, CREAM, bold=True, first=True,
     space_after=0)
add_table(s, 6.70, TOP_C + 0.60, 5.75, [
    ["What the spring does", "What the motor must do", "Result"],
    ["Lifts a little", "Carries the rest", "Better than nothing"],
    ["Lifts exactly right", "Carries nothing", "**Best case**"],
    ["Lifts too hard", "Pushes back down", "**Over-assist** — effort climbs"],
    ["Pulls the wrong way", "Carries load **and** fights spring", "**Wrong sign** — always worse"],
], col_w=[1.5, 1.85, 2.0], fs=9.5, row_h=0.40, highlight=(2,), warn=(4,),
    aligns=[PP_ALIGN.LEFT] * 3)
tf = textbox(s, 6.70, TOP_C + 2.84, 5.75, 2.4)
para(tf, "Right knees get over-assisted; left knees get wrong-signed assist, "
         "where no stiffness can help.", 10.5, CREAM, first=True, space_after=7,
     line_spacing=1.03)
para(tf, "50 of 440 knee-cells (11.4%) were wrong-signed — roughly an eighth of "
         "the sweep was uninterpretable.", 11, YELLOW, bold=True, space_after=7,
     bullet="▸", line_spacing=1.03)
para(tf, "A single rule — harmful when the assist ratio kx(θ₀−q_op)/HOLD is below "
         "0% or above 200% — predicts 92.7% of all 440 knee-cells with zero false "
         "positives.", 10.5, CREAM, space_after=7, bullet="▸", line_spacing=1.03)
para(tf, "Fix: mirror the sign per side.  Right knees θ₀ = −|θ₀|,  left knees "
         "θ₀ = +|θ₀|.", 11, YELLOW, bold=True, space_after=0, bullet="▸",
     line_spacing=1.03)
notes(s, """
[1:22-1:52] The right knees stand at plus 37 and plus 43 degrees, the left at
minus 38 and minus 41, with holding torques of opposite sign. One shared rest
angle can only help both sides at zero. Away from zero the right knees get
over-assisted - correct direction, too strong, so the motor pushes back - and
the left knees get assist in the wrong direction, where no stiffness helps. 51
of 440 knee-cells wrong-signed. The fix: give each side the sign it needs.
""")


# ================================================================ 6. MIRRORED
s = new_slide("Sweep 2 — Mirrored Rest Angle (91 runs)",
              "Same spring, same gait, sign of θ₀ matched to each leg: kx = 0.05→0.45 "
              "(9) × |θ₀| = 0°→45° (10) = 90 cells + baseline.")
px, py, pw, ph = add_pic(s, FIG / "cross_sweeps.png",
                         L, TOP_C, CW, 3.22)
add_table(s, L, py + ph + 0.30, 6.55, [
    ["Metric", "Sweep 1 — shared", "Sweep 2 — mirrored", "Change"],
    ["Best torque reduction", "34.0%", "**34.39%**", "+0.4 pts"],
    ["Wrong-sign cells", "51 / 440  (11.6%)", "**0 / 360  (0%)**", "eliminated"],
    ["Spread at own optimum", "3.96 pts", "**2.71 pts**", "1.5× tighter"],
    ["Spread at recommended", "—", "**1.15 pts**", "3.4× tighter"],
    ["kx=0.30 / θ₀=0° regression check", "34.0%", "33.7%", "within noise"],
], col_w=[2.35, 1.55, 1.6, 1.05], fs=9.5, row_h=0.325,
    aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
rect(s, 7.35, py + ph + 0.30, R - 7.35, 1.66, PERI,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
tf = textbox(s, 7.60, py + ph + 0.46, 4.85, 1.4)
para(tf, "The headline number barely moved — and that is the point.", 11, YELLOW,
     bold=True, first=True, space_after=6, line_spacing=1.03)
para(tf, "Mirroring did not buy us a bigger reduction. It bought us a sweep where "
         "all 90 cells are physically interpretable instead of an eighth being "
         "broken, and where the recommended setting helps all four knees equally.",
     10.5, CREAM, space_after=0, line_spacing=1.03)
notes(s, """
[1:52-2:14] Re-running mirrored: 34.4% versus 34.0%. Almost identical, and that
is the honest result. Mirroring did not buy a bigger number - it bought a
sweep where all 90 cells are interpretable instead of an eighth being broken,
and per-knee spread at the optimum drops from 6.1 points to 1.1.
""")


# ================================================================ 7. RIDGE
s = new_slide("A Ridge, Not a Peak — One Effective DOF",
              "Assist torque is kx · (|θ₀| + |q_op|); the two parameters multiply, so "
              "the optimum is a hyperbola.")
add_pic(s, FIG / "p2b_ridge.png", L, TOP_C, 5.35, 3.95)
tf = textbox(s, 6.30, TOP_C, R - 6.30, 3.95, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Best cell: kx = 0.15 N·m/rad, |θ₀| = ±35° → **34.39%** reduction "
         "(0.2352 → 0.1543 N·m).", 12, CREAM, first=True, space_after=9,
     bullet="▸", line_spacing=1.04)
para(tf, "Five very different parameter pairs all land within 0.8 points of each "
         "other. The design has **one effective degree of freedom, not two** — the "
         "spring only has to deliver about 0.25 N·m of assist; how that is split "
         "between stiffness and preload is a manufacturing convenience.",
     12, CREAM, space_after=9, bullet="▸", line_spacing=1.04)
para(tf, "Over-assist is the failure mode at the bottom-right: correctly signed "
         "spring, too strong, so the motor starts fighting it — down to −101% at "
         "kx=0.45 / ±45°. Interpretable, not broken.", 12, CREAM, space_after=9,
     bullet="▸", line_spacing=1.04)
para(tf, "Caveat we have to state: the ridge exits the grid at both ends. Below "
         "kx = 0.10 the best |θ₀| is 45°, the grid edge — that branch is unmapped.",
     11.5, YELLOW, space_after=0, bullet="▸", line_spacing=1.04)
add_table(s, L, 5.78, CW, [
    ["kx (N·m/rad)", "0.05", "0.10", "0.15", "0.20", "0.25", "0.30", "0.35", "0.40", "0.45"],
    ["Best |θ₀|", "±45°", "±45°", "**±35°**", "±20°", "±10°", "±5°", "±0°", "±0°", "±0°"],
    ["Reduction at that cell", "16.6%", "31.2%", "**34.4%**", "34.3%", "34.0%", "33.7%", "33.6%", "31.0%", "27.4%"],
], col_w=[2.0] + [1.0] * 9, fs=9.5, hdr_fs=9.5, row_h=0.315,
    aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 9)
notes(s, """
[2:14-2:45] The most useful structural result. Assist torque is stiffness times
angle, so the two parameters trade off - the optimum is a hyperbolic ridge, not
a point. As stiffness rises the best rest angle falls, and everything from 0.15
to 0.35 delivers 33.6 to 34.4 percent. So the defensible claim is not "the
optimum is 0.15 and 35 degrees" - it is that a correctly sized assist of about
0.25 newton-metres removes a third of knee effort, and this whole locus
delivers it. Good for manufacturing tolerance.
""")


# ================================================================ 8. PER-KNEE
s = new_slide("Per-Knee Symmetry — What Mirroring Bought",
              "Spread = best knee minus worst knee, in percentage points of torque "
              "reduction.")
add_pic(s, FIG / "cross_asymmetry.png", L, TOP_C, 5.62, 3.55)
rect(s, L, 5.30, 5.62, 1.58, PERI, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.10)
tf = textbox(s, L + 0.28, 5.46, 5.06, 1.3)
para(tf, "Why this matters for hardware: at Sweep 1's optimum the knees span "
         "31.5–35.4%, so **one spring cannot serve every leg equally**. At "
         "1.15 points, it can.", 11, CREAM, first=True, space_after=0,
     bullet="▸", line_spacing=1.04)
add_table(s, 6.45, TOP_C + 0.22, R - 6.45, [
    ["Configuration", "FR", "BR", "BL", "FL", "Spread"],
    ["Sweep 1 optimum  (0.30, 0°)", "33.6", "35.2", "35.4", "31.5", "3.96 pts"],
    ["Sweep 2 recommended  (0.20, ±15°)", "33.9", "34.5", "34.6", "33.4", "**1.1 pts**"],
    ["Sweep 2 torque optimum  (0.15, ±35°)", "34.3", "34.6", "35.6", "32.9", "2.7 pts"],
    ["Over-assist corner  (0.45, ±45°)", "−99.5", "−97.9", "−95.3", "−110.9", "15.6 pts"],
], col_w=[2.6, 0.62, 0.62, 0.62, 0.62, 0.85], fs=9.5, row_h=0.40, highlight=(2,),
    warn=(4,), aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 5)
tf = textbox(s, 6.45, TOP_C + 2.35, R - 6.45, 3.0)
para(tf, "Across the whole grid the spread runs from 1.1 to 15.6 points, median 8.9. "
         "The recommended cell is the **most bilaterally symmetric of all 90** — it "
         "happens to sit near all four knees' individual optima at once.",
     11, CREAM, first=True, space_after=9, line_spacing=1.04)
para(tf, "**0 wrong-sign cells in 360 checked**, versus 50 of 440 before mirroring.",
     11.5, YELLOW, space_after=9, bullet="▸", line_spacing=1.04)
para(tf, "FL is consistently the weakest knee and BL the strongest. That is not a "
         "mirroring failure — the four legs have genuinely different measured "
         "holding torques (0.246 → 0.264 N·m) and stance angles. Only per-knee "
         "stiffness would remove the last point of it.", 11, CREAM, space_after=0,
     bullet="▸", line_spacing=1.04)
notes(s, """
[2:45-3:05] Why symmetry matters for hardware: at Sweep 1's optimum the four
knees span 31.5 to 35.4 percent, so one spring cannot serve every leg equally.
That is a 3.96 point spread; the mirrored recommendation is 1.15, tightest of
all ninety cells. The residual is real, not a bug - the four legs measure
slightly different holding torques.
""")


# ================================================================ 9. CoT
s = new_slide("Cost of Transport — Three Definitions",
              "CoT = E / (m·g·d),  m = 1.398 kg, g = 9.8 m/s², all 12 joints. "
              "The three variants differ only in how the energy E is counted.")
px, py, pw, ph = add_pic(s, FIG / "p2b_cot_bars.png",
                         L, TOP_C, 11.15, 3.62)
add_table(s, L, py + ph + 0.30, 7.35, [
    ["Energy definition", "Baseline", "Best cell", "Best value", "Change", "r vs torque red."],
    ["Mechanical   Σ|τ·dθ|", "2.7149", "0.25 / ±15°", "2.3012", "−15.2%", "−0.787"],
    ["Positive work   Σmax(0, τ·dθ)", "2.1830", "0.25 / ±15°", "1.8178", "−16.7%", "−0.366"],
    ["Electrical proxy   ∫τ²dt", "0.8779", "0.20 / ±15°", "**0.5734**", "**−34.7%**", "**−0.986**"],
], col_w=[2.35, 0.85, 1.0, 0.95, 0.85, 1.15], fs=9.5, row_h=0.34, highlight=(3,),
    aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 5)
rect(s, 8.05, py + ph + 0.30, R - 8.05, 1.40, PERI,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
tf = textbox(s, 8.28, py + ph + 0.42, 4.2, 1.2)
para(tf, "A motor holding a static load burns current but does zero mechanical "
         "work. The spring's whole job is to cancel static holding torque — so "
         "**mechanical CoT is measuring the wrong thing** and understates the "
         "benefit ~5×. The I²R proxy tracks the torque saving at r = −0.986.",
     10, CREAM, first=True, space_after=0, line_spacing=1.03)
notes(s, """
[3:05-3:40] Cost of transport, over all twelve joints and the whole robot.
Mechanical work says the spring saves only 7% - that looks like a failure, but
it is an artifact: mechanical work is torque times angle change, so a motor
holding a static load registers zero work while still burning current, and
cancelling static holding torque is the spring's entire job. The electrical
proxy, integral of torque squared, is the copper-loss surrogate - it drops
34.7%, tracking the torque reduction at r equals minus 0.986. That is the
efficiency number we quote. It is a proxy: joules would need the servo's torque
constant and resistance.
""")


# ================================================================ 10. METRICS
s = new_slide("Metric by Metric — What Moves, What Doesn't",
              "Left: baseline vs the torque-optimal cell, with each metric's "
              "correlation to torque reduction across all 90 cells.")
add_table(s, L, TOP_C, 6.55, [
    ["Metric", "Baseline", "At optimum\n(0.15, ±35°)", "Change", "r vs red.", "Independent?"],
    ["Mean knee effort (N·m)", "0.2352", "**0.1543**", "**−34.4%**", "—", "headline"],
    ["RMS effort (N·m)", "0.2863", "0.2126", "−25.7%", "−0.996", "collinear"],
    ["p99 demand (N·m)", "0.9311", "0.8745", "−6.1%", "−0.852", "partly"],
    ["Peak demand (N·m)", "0.9831", "0.9502", "−3.3%", "—", "partly"],
    ["Torque variance", "0.0266", "0.0214", "−19.5%", "−0.596", "partly"],
    ["Saturation (%)", "0.69", "0.75", "+0.06 pts", "−0.457", "**yes**"],
    ["Mean tracking error", "4.299°", "3.395°", "−21.0%", "−0.996", "collinear"],
    ["Forward displacement (m)", "0.3256", "0.3297", "+1.3%", "+0.192", "**yes**"],
], col_w=[2.0, 0.85, 1.15, 0.9, 0.8, 1.0], fs=9.5, hdr_fs=8.5, row_h=0.335,
    hdr_h=0.52, highlight=(1,),
    aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 5)
tf = textbox(s, L, TOP_C + 3.52, 6.55, 1.5)
para(tf, "RMS falls less than mean (−25.7% vs −34.4%): the spring removes the DC "
         "gravity bias, not the AC dynamic component. The continuous thermal "
         "benefit is real but smaller than the headline.", 10.5, CREAM, first=True,
     space_after=6, bullet="▸", line_spacing=1.03)
para(tf, "Tracking error is collinear at r = −0.996 — it restates the effort "
         "result and **cannot** be used as independent evidence that gait quality "
         "survived. Saturation and displacement are the only genuinely independent "
         "checks, and both stay healthy.", 10.5, YELLOW, space_after=0, bullet="▸",
     line_spacing=1.03)

rect(s, 7.42, TOP_C, R - 7.42, 5.15, PERI, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.05)
tf = textbox(s, 7.68, TOP_C + 0.20, 4.8, 0.62)
para(tf, "RECOMMENDED CONFIGURATION", 11.5, YELLOW, bold=True, first=True,
     space_after=1)
para(tf, "kx = 0.20 N·m/rad,  |θ₀| = ±15°", 13.5, CREAM, font=SERIF, bold=True,
     space_after=0)
add_table(s, 7.68, TOP_C + 0.94, 4.8, [
    ["Property", "Value", "Rank of 90"],
    ["Electrical CoT", "0.5734", "**#1**"],
    ["RMS effort", "0.2121 N·m", "**#1**"],
    ["Per-knee spread", "1.1 pts", "**#1**"],
    ["Actuator saturation", "0.31%", "15–24 of 90"],
    ["Torque reduction", "34.12%", "0.27 pts off #1"],
    ["p99 demand", "0.8636 N·m", "8% under rating"],
    ["Mechanical CoT", "2.4826", "−8.6% vs base"],
    ["Tracking error", "3.54°", "−17.6%"],
    ["Displacement", "0.3309 m", "at grid mean"],
], col_w=[1.55, 1.25, 1.5], fs=9.5, row_h=0.335,
    aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
tf = textbox(s, 7.68, TOP_C + 4.44, 4.8, 0.55)
para(tf, "Not the winner on torque reduction — but **#1 of 90** on electrical "
         "cost, RMS effort and bilateral symmetry.", 10, CREAM, italic=True,
     first=True, space_after=0, line_spacing=1.02)
notes(s, """
[3:40-4:12] Two honest flags. RMS effort falls 26% while mean falls 34% - the
spring cancels the constant gravity bias, not the dynamic component, so the
thermal benefit is smaller than the headline. And tracking error correlates at
minus 0.996 with effort, so it is the same measurement restated - not
independent proof the gait survived. Saturation and displacement are the only
independent checks, and both are fine. On the right: our recommended
configuration is not the torque winner - it is first of ninety on electrical
cost, RMS, symmetry and saturation.
""")


# ================================================================ 11. FREQ
s = new_slide("Speed Lever 1 — Replay Frequency",
              "The same 16 waypoints, played faster or slower. Per-step jump is 2.989° "
              "in all three runs — so only speed changed.")
px, py, pw, ph = add_pic(s, FIG / "p3a_bars.png",
                         L, TOP_C, CW, 2.92)
add_table(s, L, py + ph + 0.28, CW, [
    ["target_freq", "Cycle time", "Mean effort", "RMS effort", "Peak demand", "Saturation", "Tracking err (mean/RMS/peak)"],
    ["5 Hz  (slower)", "3.2 s", "0.2101  (−8.3%)", "0.2535  (−9.6%)", "0.953  (−1.9%)", "**0.19%  (¼×)**", "3.76° / 6.58° / 21.0°"],
    ["10 Hz  (baseline)", "1.6 s", "0.2292", "0.2804", "0.972", "0.75%", "4.11° / 6.81° / 20.8°"],
    ["20 Hz  (faster)", "0.8 s", "0.2424  (+5.8%)", "0.3265  (+16.5%)", "**1.235  (+27.1%)**", "**4.88%  (6.5×)**", "4.19° / 6.92° / 21.2°"],
], col_w=[1.5, 1.0, 1.45, 1.45, 1.45, 1.35, 2.1], fs=9.5, hdr_fs=9, row_h=0.345,
    highlight=(1,), warn=(3,),
    aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 6)
tf = textbox(s, L, py + ph + 1.72, CW, 0.62)
para(tf, "Speeding up is expensive at the **peak**, not at the mean: the same 2.989° "
         "jump in half the time makes the D-term react harder. Mean effort barely "
         "moves because it is set by the unchanged stance holding torque.",
     11, CREAM, first=True, space_after=0, bullet="▸", line_spacing=1.04)
notes(s, """
[4:12-4:40] Two speed levers, both with no spring. First, replay frequency -
same sixteen waypoints, faster or slower; the command log confirms the per-step
jump is 2.989 degrees in all three, so only speed changed. 20 Hz costs 27% more
peak torque and multiplies saturation by six and a half; 5 Hz drops saturation
to a quarter. Mean effort barely moves - it is set by stance holding torque.
""")


# ================================================================ 12. STEPS
s = new_slide("Speed Lever 2 — Trajectory Resolution",
              "Sampling the same Bézier swing arc with more or fewer waypoints, at a "
              "fixed 10 Hz. This is the lever that actually fixes peak torque.")
px, py, pw, ph = add_pic(s, FIG / "p3b_bars.png",
                         L, TOP_C, CW, 2.88)
add_table(s, L, py + ph + 0.26, 7.6, [
    ["NUM_DATA_POINTS", "Cycle", "Mean effort", "RMS", "Peak demand", "Sat.", "Track err", "Step jump"],
    ["8   ⚠ degenerate", "0.8 s", "0.1749", "0.2090", "0.689", "0.00%", "1.00°", "0.483°"],
    ["16   (baseline)", "1.6 s", "0.2246", "0.2756", "0.965", "0.63%", "4.05°", "2.989°"],
    ["32   (finer)", "3.2 s", "0.1994", "0.2419", "**0.495  (−48.7%)**", "**0.00%**", "2.92°", "1.610°"],
], col_w=[1.55, 0.7, 1.0, 0.85, 1.45, 0.66, 0.83, 0.85], fs=9, hdr_fs=8.5,
    row_h=0.345, highlight=(3,), warn=(1,),
    aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 7)
tf = textbox(s, L, py + ph + 1.80, 7.6, 1.05)
para(tf, "**The N = 8 trap.**  n_swing = int(8 × 0.25) = 2, so linspace samples only "
         "t = 0 and t = 1 — both already at stance height. The foot never lifts; it "
         "slides. Its \"better\" numbers are a drag, not a step. Hard cliff: any lift "
         "needs 3 swing samples, i.e. **N ≥ 12**. Excluded from all conclusions.",
     10.5, YELLOW, first=True, space_after=0, bullet="⚠", line_spacing=1.03)

rect(s, 8.42, py + ph + 0.26, R - 8.42, 2.05, PERI,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.045)
tf = textbox(s, 8.66, py + ph + 0.42, 3.85, 2.0)
para(tf, "MATCHED AT ≈3.2 s CYCLE TIME", 10, YELLOW, bold=True, first=True,
     space_after=5)
para(tf, "5 Hz replay → peak demand **0.953**", 10.5, CREAM, space_after=3)
para(tf, "32 waypoints → peak demand **0.495**", 10.5, CREAM, space_after=7)
para(tf, "At the same cycle time, finer sampling halves peak torque where slower "
         "replay does nothing. Slowing down gives the controller more time to chase "
         "the same discontinuity; finer sampling shrinks the discontinuity itself.",
     10, CREAM, space_after=0, line_spacing=1.03)
notes(s, """
[4:40-5:20] Second lever, resolution - same speed, more waypoints on the same
curve. A trap first: at eight points the swing phase gets two samples, both
endpoints, both at stance height, so the foot never lifts, it slides. Its
numbers look good because it is doing less, so we excluded it - the cliff is
exact at twelve points. The valid comparison is 16 versus 32, and it is the
strongest result here: peak demand down 48.7%, saturation to zero, with no
spring. And at matched cycle time, finer sampling halves peak torque while
slower replay does nothing - slowing down only gives the controller longer to
chase the same discontinuity; finer sampling removes it.
""")


# ================================================================ 13. CLOSE
s = new_slide("Where We Stand", None)
rect(s, L, TOP_N, 6.05, 2.96, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
tf = textbox(s, L + 0.26, TOP_N + 0.20, 5.55, 3.0)
para(tf, "SUPPORTABLE CLAIMS", 11, INK, bold=True, first=True, space_after=7)
for t in [
    "A mirrored parallel knee spring removes **~34% of mean knee motor torque** "
    "and ~26% of RMS effort — 91 runs, combined across four knees.",
    "Electrical (copper-loss) cost of transport falls **~35%**, 0.878 → 0.573.",
    "Mirroring gives **0 wrong-sign cells in 360** (was 50/440) and cuts bilateral "
    "asymmetry from 3.96 to **1.15 points**.",
    "The design has **one effective DOF** — a locus of (kx, θ₀) all deliver ~34%.",
    "Finer trajectory sampling cuts peak demand **−49%**, independent of the spring.",
]:
    para(tf, t, 10, INK2, space_after=5, bullet="✓", line_spacing=1.02)

rect(s, 6.92, TOP_N, R - 6.92, 2.96, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
     radius=0.05)
tf = textbox(s, 7.18, TOP_N + 0.20, 5.3, 3.0)
para(tf, "WHAT WE DO NOT CLAIM", 11, INK, bold=True, first=True, space_after=7)
for t in [
    "Mechanical work does **not** drop in proportion to torque — only 6.2% against "
    "34.4%. Any efficiency claim from mechanical CoT alone understates this device ~5×.",
    "Electrical CoT **in joules** — the proxy needs the servo's kt and R first.",
    "Tracking error as independent proof of gait quality — collinear at r = −0.996.",
    "A smaller motor. Peak demand is control-limited, not spring-limited: baseline "
    "p99 is 0.931 N·m against a 0.941 N·m rating, and no spring setting fixes that.",
    "n = 1 per cell — differences under ~1 point are not statistically separable.",
]:
    para(tf, t, 10, INK2, space_after=5, bullet="✗", line_spacing=1.02)

rect(s, L, 4.42, CW, 2.46, PERI, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055)
tf = textbox(s, L + 0.28, 4.62, 5.35, 2.1)
para(tf, "LITERATURE CONTEXT", 10.5, YELLOW, bold=True, first=True, space_after=5)
para(tf, "Published passive-compliance savings span **17–50%**: compliant feet ~17%, "
         "hexapod gait + torque optimisation 22–39%, nonlinear elastic joints up to "
         "50%, trajectory-optimised biped coupling >50%.", 10, CREAM, space_after=5,
     line_spacing=1.03)
para(tf, "Our ~34% sits solidly in that range using the simplest possible hardware — "
         "a fixed linear spring, no clutch, no adaptive mechanism.", 10, CREAM,
     space_after=0, line_spacing=1.03)
tf = textbox(s, 6.55, 4.62, 5.95, 2.1)
para(tf, "NEXT STEPS", 10.5, YELLOW, bold=True, first=True, space_after=5)
for t in [
    "Three repeats at the ridge to make sub-1-point differences separable.",
    "Extend the grid to |θ₀| > 45° at low stiffness — the ridge exits the box there.",
    "Raise NUM_DATA_POINTS 16 → 80 to remove the 10 Hz step discontinuity, then "
    "re-measure peak torque.",
    "Add an explicit fall / body-tilt gate instead of inferring success from "
    "displacement.",
    "Get kt and R from the servo datasheet to convert the electrical proxy to joules.",
]:
    para(tf, t, 10, CREAM, space_after=4, bullet="→", line_spacing=1.02)
notes(s, """
[5:20-6:00] What we stand behind: 34% mean knee torque, 35% electrical cost of
transport, zero wrong-sign cells, 1.1 point symmetry, one-degree-of-freedom
design surface. What we do not claim, before being asked: mechanical work only
drops 6%, so no efficiency claim from it; the electrical figure is a proxy, not
joules; tracking error is not independent evidence; and this does not justify a
smaller motor, because peak demand is limited by the 10 Hz stepped setpoint.
Our 34% sits inside the published 17-to-50% range, with the simplest hardware
in it. Next: repeats at the ridge, extend the grid, raise trajectory
resolution, and get the motor constants for real joules.
""")


prs.save(str(OUT))
print(f"Wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
